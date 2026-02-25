import os
from pptx import Presentation

ppt_dir = "/root/eco_design/ecofinal/ecofinal/金融数据分析与智能量化交易应用PPT_2025"
md_dir = "/root/eco_design/ecofinal/ecofinal/金融数据分析与智能量化交易应用PPT_2025_MD"

os.makedirs(md_dir, exist_ok=True)

for filename in sorted(os.listdir(ppt_dir)):
    if filename.endswith(".pptx"):
        ppt_path = os.path.join(ppt_dir, filename)
        md_filename = filename.replace(".pptx", ".md")
        md_path = os.path.join(md_dir, md_filename)
        
        print(f"Converting {filename} to Markdown...")
        try:
            prs = Presentation(ppt_path)
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(f"# {filename}\n\n")
                for i, slide in enumerate(prs.slides):
                    f.write(f"## Slide {i+1}\n\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                f.write(f"{text}\n\n")
            print(f"  -> Saved to {md_filename}")
        except Exception as e:
            print(f"  -> Error converting {filename}: {e}")
