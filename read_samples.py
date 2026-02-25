import os
from pptx import Presentation
import nbformat

def read_ppt(file_path):
    print(f"\n--- Reading PPT: {os.path.basename(file_path)} ---")
    try:
        prs = Presentation(file_path)
        print(f"Total slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides[:3]): # Read first 3 slides
            print(f"  Slide {i+1}:")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(f"    - {shape.text[:50]}...") # Print first 50 chars
    except Exception as e:
        print(f"Error reading PPT: {e}")

def read_ipynb(file_path):
    print(f"\n--- Reading Notebook: {os.path.basename(file_path)} ---")
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            nb = nbformat.read(f, as_version=4)
            print(f"Total cells: {len(nb.cells)}")
            code_cells = [c for c in nb.cells if c.cell_type == 'code']
            print(f"Code cells: {len(code_cells)}")
            if code_cells:
                print("First code cell snippet:")
                print(f"  {code_cells[0].source[:100]}...")
    except Exception as e:
        print(f"Error reading Notebook: {e}")

ppt_dir = "/root/eco_design/ecofinal/ecofinal/金融数据分析与智能量化交易应用PPT_2025"
code_dir = "/root/eco_design/ecofinal/ecofinal/金融数据分析与智能量化交易应用代码_2025"

# Read a sample PPT
ppt_files = [f for f in os.listdir(ppt_dir) if f.endswith('.pptx')]
if ppt_files:
    read_ppt(os.path.join(ppt_dir, ppt_files[0]))

# Read a sample Notebook
ipynb_files = [f for f in os.listdir(code_dir) if f.endswith('.ipynb')]
if ipynb_files:
    read_ipynb(os.path.join(code_dir, ipynb_files[0]))
